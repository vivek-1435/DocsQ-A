import os
import warnings
from pathlib import Path

# suppress langchain deprecation warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)

from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser


SYSTEM_PROMPT = """\
You are an expert, highly intelligent AI assistant analyzing the provided document context.

Instructions:
1. Answer the user's question by thoroughly analyzing the provided document context.
2. Ground all your answers strictly in the facts and details present in the context. Do not invent details or assume facts that are not present.
3. If the user asks an indirectly connected, analytical, synthesis, or subjective question (such as comparing options, evaluating key aspects, identifying the "most prominent" or "best" items, or summarizing complex themes), provide a helpful, professional, and synthesized analysis based *only* on the facts in the document. Note clearly if the document does not explicitly specify a ranking or direct answer, but analyze and compare the details actually present in the context to provide a useful, well-reasoned answer.
4. If the question is about a completely unrelated topic, general knowledge, or external subjects that have absolutely nothing to do with the document context (e.g., "What is the capital of France?", "tell me a recipe"), state clearly:
"I couldn't find relevant information in the document to answer this question."

Context:
{context}
"""


def _format_docs(docs: list) -> str:
    return "\n\n---\n\n".join(doc.page_content for doc in docs)


class ParallelGoogleGenerativeAIEmbeddings(GoogleGenerativeAIEmbeddings):
    def embed_documents(
        self,
        texts: list[str],
        *,
        batch_size: int = 100,
        task_type: str | None = None,
        titles: list[str] | None = None,
        output_dimensionality: int | None = None,
    ) -> list[list[float]]:
        import time
        import re
        import random
        from concurrent.futures import ThreadPoolExecutor

        # Calculate dynamic batch size targeting max 150,000 characters per batch
        if texts:
            chunk_len = len(texts[0])
            dynamic_batch_size = max(1, min(100, int(150000 / chunk_len)))
        else:
            dynamic_batch_size = 100

        print(f"📊 Dynamic batch size: first chunk length is {chunk_len if texts else 0} chars. Using batch_size={dynamic_batch_size}")
        batches = [texts[i:i + dynamic_batch_size] for i in range(0, len(texts), dynamic_batch_size)]
        
        def _embed_batch(batch_texts):
            retries = 6
            for attempt in range(retries):
                try:
                    return super(ParallelGoogleGenerativeAIEmbeddings, self).embed_documents(
                        batch_texts,
                        batch_size=len(batch_texts),
                        task_type=task_type,
                        output_dimensionality=output_dimensionality
                    )
                except Exception as e:
                    err_str = str(e)
                    is_rate_limit = False
                    if hasattr(e, 'status_code') and e.status_code == 429:
                        is_rate_limit = True
                    elif "429" in err_str or "RESOURCE_EXHAUSTED" in err_str:
                        is_rate_limit = True
                    
                    if is_rate_limit and attempt < retries - 1:
                        sleep_time = (2.5 ** attempt) + 1.0 + random.uniform(0.5, 3.5)
                        match = re.search(r"retry in ([\d\.]+)s", err_str)
                        if match:
                            sleep_time = float(match.group(1)) + 1.0 + random.uniform(0.5, 2.5)
                        print(f"\n⚠️ Rate limit (429) hit. Attempt {attempt+1}/{retries}. Sleeping {sleep_time:.2f}s...")
                        time.sleep(sleep_time)
                    else:
                        raise e
            
        results = []
        max_workers = 5
        print(f"Embedding {len(texts)} chunks in parallel with {max_workers} threads...")
        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            futures = []
            for batch in batches:
                futures.append(executor.submit(_embed_batch, batch))
                # Add a 0.5s pacing delay to prevent rate limit spikes
                time.sleep(0.5)
                
            for future in futures:
                results.extend(future.result())
                
        return results


class RAGPipeline:
    def __init__(self):
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            raise ValueError("GEMINI_API_KEY not found in environment variables")

        self.embeddings = ParallelGoogleGenerativeAIEmbeddings(
            model="models/gemini-embedding-001",
            google_api_key=api_key,
        )
        self.llm = ChatGoogleGenerativeAI(
            model="gemini-2.5-flash-lite",
            google_api_key=api_key,
            temperature=0.3,
        )
        self.vector_store = None
        self.retriever = None
        self.doc_name = ""

    def _load_document(self, file_path: str) -> list:
        from langchain_core.documents import Document
        ext = Path(file_path).suffix.lower()

        if ext == ".pdf":
            from langchain_community.document_loaders import PyPDFLoader
            return PyPDFLoader(file_path).load()

        elif ext == ".docx":
            from langchain_community.document_loaders import Docx2txtLoader
            return Docx2txtLoader(file_path).load()

        elif ext == ".txt":
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="latin-1") as f:
                    content = f.read()
            return [Document(page_content=content, metadata={"source": Path(file_path).name})]

        elif ext == ".csv":
            import csv
            text_lines = []
            try:
                with open(file_path, "r", encoding="utf-8") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        row_vals = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                        if row_vals:
                            text_lines.append(" | ".join(row_vals))
            except UnicodeDecodeError:
                with open(file_path, "r", encoding="latin-1") as f:
                    reader = csv.reader(f)
                    for row in reader:
                        row_vals = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                        if row_vals:
                            text_lines.append(" | ".join(row_vals))
            content = "\n".join(text_lines)
            return [Document(page_content=content, metadata={"source": Path(file_path).name})]

        elif ext == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            text_lines = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_lines.append(f"--- Sheet: {sheet_name} ---")
                for row in sheet.iter_rows(values_only=True):
                    row_vals = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                    if row_vals:
                        text_lines.append(" | ".join(row_vals))
            content = "\n".join(text_lines)
            return [Document(page_content=content, metadata={"source": Path(file_path).name})]

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            text_lines = []
            for i, slide in enumerate(prs.slides):
                text_lines.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    try:
                        if hasattr(shape, "text") and shape.text.strip():
                            text_lines.append(shape.text.strip())
                        if hasattr(shape, "has_table") and shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                                if row_vals:
                                    text_lines.append(" | ".join(row_vals))
                    except Exception as shape_err:
                        # Defensive: Skip any corrupted or non-standard PPTX vector shapes
                        print(f"Skipping shape in PPTX slide due to error: {shape_err}")
                        continue
            content = "\n".join(text_lines)
            return [Document(page_content=content, metadata={"source": Path(file_path).name})]

        raise ValueError(f"Unsupported file type: {ext}")

    def process_document(self, file_path: str, original_name: str = ""):
        self.doc_name = original_name or Path(file_path).name

        documents = self._load_document(file_path)
        
        # Check if we successfully extracted any selectable text
        total_text = ""
        if documents:
            total_text = "".join(doc.page_content for doc in documents).strip()

        # If text is empty or extremely short, trigger Gemini Multimodal OCR fallback for PDFs
        if len(total_text) < 15:
            ext = Path(file_path).suffix.lower()
            if ext == ".pdf":
                print(f"⚠️ Selectable text is empty for {self.doc_name}. Falling back to Gemini Multimodal OCR...")
                try:
                    from google import genai
                    from google.genai import types
                    
                    api_key = os.getenv("GEMINI_API_KEY")
                    client = genai.Client(api_key=api_key)
                    
                    with open(file_path, "rb") as f:
                        file_bytes = f.read()
                        
                    response = client.models.generate_content(
                        model='gemini-2.5-flash-lite',
                        contents=[
                            types.Part.from_bytes(
                                data=file_bytes,
                                mime_type="application/pdf"
                            ),
                            "Transcribe all text from this scanned document. Please preserve layout, headers, tables, and write out all text exactly as it appears. Do not summarize or add introduction."
                        ]
                    )
                    
                    ocr_text = response.text
                    if ocr_text and ocr_text.strip():
                        total_text = ocr_text.strip()
                except Exception as ocr_err:
                    print(f"⚠️ Gemini OCR fallback failed: {ocr_err}")

        # Final foolproof check: if no content could be extracted, use a placeholder
        if not total_text:
            total_text = f"This document '{self.doc_name}' contains no readable text or is empty."

        # Truncate content to protect API limit and ensure fast response times
        MAX_CHARS = 1500000
        if len(total_text) > MAX_CHARS:
            print(f"⚠️ Document content length ({len(total_text)}) exceeds MAX_CHARS ({MAX_CHARS}). Truncating to optimize speed.")
            total_text = total_text[:MAX_CHARS]

        # Re-wrap content in a single Document object for RecursiveCharacterTextSplitter
        from langchain_core.documents import Document
        documents = [Document(page_content=total_text, metadata={"source": self.doc_name})]
        total_chars = len(total_text)

        # Adaptive chunking thresholds capped at 8000 for gemini-embedding-001 limits
        if total_chars < 500 * 1024: # Under 500 KB
            chunk_size = 1500
            chunk_overlap = 150
        elif total_chars < 2 * 1024 * 1024: # 500 KB to 2 MB
            chunk_size = 4000
            chunk_overlap = 400
        else: # Over 2 MB (capped at 8000 chars)
            chunk_size = 8000
            chunk_overlap = 800
            
        print(f"📊 Adaptive chunking: document character count is {total_chars}. Using chunk_size={chunk_size}, chunk_overlap={chunk_overlap}")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        chunks = splitter.split_documents(documents)
        
        # If chunks is empty, use the placeholder text
        if not chunks:
            from langchain_core.documents import Document
            placeholder_text = f"This document '{self.doc_name}' contains no readable text or is empty."
            chunks = [Document(page_content=placeholder_text, metadata={"source": self.doc_name})]

        self.vector_store = FAISS.from_documents(chunks, self.embeddings)
        self.retriever = self.vector_store.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 6},
        )

    def save_index(self, folder_path: str):
        if self.vector_store:
            Path(folder_path).mkdir(parents=True, exist_ok=True)
            self.vector_store.save_local(folder_path)

    def load_index(self, folder_path: str, doc_name: str = ""):
        self.doc_name = doc_name
        self.vector_store = FAISS.load_local(
            folder_path,
            self.embeddings,
            allow_dangerous_deserialization=True
        )
        self.retriever = self.vector_store.as_retriever(
            search_type="similarity",
            search_kwargs={"k": 6},
        )

    def ask(self, question: str) -> dict:
        if not self.retriever:
            raise ValueError("No document has been processed yet.")

        # 1. Retrieve the source documents once (cuts embedding search latency in half!)
        source_docs = self.retriever.invoke(question)

        # 2. Format the context for the prompt
        context_str = _format_docs(source_docs)

        # 3. Create the chat prompt template and run the reasoning chain
        prompt = ChatPromptTemplate.from_messages(
            [
                ("system", SYSTEM_PROMPT),
                ("human", "{question}"),
            ]
        )
        llm_chain = prompt | self.llm | StrOutputParser()
        answer = llm_chain.invoke({"context": context_str, "question": question})

        # build unique list of citations
        seen: set[str] = set()
        sources = []
        for doc in source_docs:
            snippet = doc.page_content.strip()
            if snippet in seen:
                continue
            seen.add(snippet)
            sources.append(
                {
                    "content": snippet[:300] + ("…" if len(snippet) > 300 else ""),
                    "page": doc.metadata.get("page"),
                    "source": doc.metadata.get("source", self.doc_name),
                }
            )

        return {
            "answer": answer,
            "sources": sources,
        }

