import csv
import io
import logging
from io import BytesIO
from pathlib import Path

from langchain_core.documents import Document

log = logging.getLogger(__name__)

ZIP_EXTENSIONS = {".docx", ".xlsx", ".pptx"}


def _check_magic_bytes(data: bytes, ext: str) -> None:
    if not data:
        raise ValueError("File is empty.")
    if ext == ".pdf" and not data[:4].startswith(b"%PDF"):
        raise ValueError("Not a valid PDF.")
    if ext in ZIP_EXTENSIONS and not data[:4].startswith(b"PK\x03\x04"):
        raise ValueError(f"Not a valid {ext.upper()} file.")


def _doc(filename: str, text: str) -> Document:
    return Document(page_content=text, metadata={"source": filename})


def _decode(raw: bytes) -> str:
    try:
        return raw.decode("utf-8")
    except UnicodeDecodeError:
        return raw.decode("latin-1")


def _row_text(row) -> str:
    return " | ".join(str(c).strip() for c in row if c is not None and str(c).strip())


def load_document(file_bytes: bytes, filename: str) -> list[Document]:
    ext = Path(filename).suffix.lower()
    _check_magic_bytes(file_bytes, ext)

    if ext == ".pdf":
        from pypdf import PdfReader
        reader = PdfReader(BytesIO(file_bytes))
        return [
            Document(page_content=page.extract_text() or "", metadata={"source": filename, "page": i})
            for i, page in enumerate(reader.pages)
        ]

    if ext == ".docx":
        from docx import Document as Docx
        doc = Docx(BytesIO(file_bytes))
        return [_doc(filename, "\n".join(p.text for p in doc.paragraphs))]

    if ext == ".txt":
        return [_doc(filename, _decode(file_bytes))]

    if ext == ".csv":
        lines = []
        for enc in ("utf-8", "latin-1"):
            try:
                reader = csv.reader(io.TextIOWrapper(BytesIO(file_bytes), encoding=enc))
                lines = [_row_text(row) for row in reader if _row_text(row)]
                break
            except UnicodeDecodeError:
                continue
        return [_doc(filename, "\n".join(lines))]

    if ext == ".xlsx":
        import openpyxl
        wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True)
        lines = []
        for name in wb.sheetnames:
            lines.append(f"--- Sheet: {name} ---")
            lines += [_row_text(r) for r in wb[name].iter_rows(values_only=True) if _row_text(r)]
        return [_doc(filename, "\n".join(lines))]

    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(BytesIO(file_bytes))
        lines = []
        for i, slide in enumerate(prs.slides):
            lines.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
                    if getattr(shape, "has_table", False):
                        lines += [_row_text(c.text for c in row.cells) for row in shape.table.rows]
                except Exception as err:
                    log.warning("Skipped PPTX shape: %s", err)
        return [_doc(filename, "\n".join(lines))]

    raise ValueError(f"Unsupported file type: {ext}")
